"""
Presentation Generator — Bluestock Mutual Fund Analytics
Generates 12-slide PowerPoint presentation
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd

BASE_DIR = Path(__file__).parent.parent
REPORTS_DIR = BASE_DIR / 'reports'
CHARTS_DIR = REPORTS_DIR / 'charts'
PROCESSED_DIR = BASE_DIR / 'data' / 'processed'

PURPLE = RGBColor(0x5B, 0x4F, 0xCF)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1E, 0x2E)

def add_text(slide, text, left, top, width, height, fontsize=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

def set_bg(slide, prs, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def generate_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Slide 1 — Title
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs, DARK)
    add_text(slide, "BLUESTOCK FINTECH", 1, 1.5, 11, 1, fontsize=36, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(slide, "Mutual Fund Analytics Platform", 1, 2.8, 11, 1, fontsize=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "End-to-End Data Engineering, ETL Pipeline & Interactive Dashboard", 1, 3.8, 11, 0.8, fontsize=14, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(slide, "Data Analyst Intern  |  July 2026", 1, 5.5, 11, 0.5, fontsize=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Slide 2 — Problem & Objective
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs, WHITE)
    add_text(slide, "Problem & Objective", 0.5, 0.3, 12, 0.8, fontsize=28, bold=True, color=PURPLE)
    problems = [
        "❌  NAV, AUM, SIP data scattered across multiple sources in different formats",
        "❌  No unified platform to compare fund performance on a risk-adjusted basis",
        "❌  Investors cannot track whether funds outperform their benchmark index",
        "❌  Monthly reports are static PDFs — no interactivity or drill-down capability",
    ]
    objectives = [
        "✅  Build ETL pipeline consolidating all data into a single SQLite database",
        "✅  Compute risk-return metrics: CAGR, Sharpe, Alpha, Beta, VaR, Max Drawdown",
        "✅  Compare fund returns vs Nifty 50 and Nifty 100 benchmarks",
        "✅  Deliver interactive Power BI dashboard with 4 analytical pages",
    ]
    for i, p in enumerate(problems):
        add_text(slide, p, 0.5, 1.2 + i*0.55, 6, 0.5, fontsize=11, color=DARK)
    for i, o in enumerate(objectives):
        add_text(slide, o, 6.8, 1.2 + i*0.55, 6, 0.5, fontsize=11, color=DARK)
    add_text(slide, "Problems", 0.5, 0.9, 5, 0.4, fontsize=13, bold=True, color=ORANGE)
    add_text(slide, "Solutions", 6.8, 0.9, 5, 0.4, fontsize=13, bold=True, color=PURPLE)

    # Slide 3 — Data Sources
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs, WHITE)
    add_text(slide, "Data Sources & Datasets", 0.5, 0.3, 12, 0.8, fontsize=28, bold=True, color=PURPLE)
    datasets = [
        ("40 Schemes", "fund_master.csv — AMFI codes, categories, expense ratios"),
        ("46,000 rows", "nav_history.csv — Daily NAV 2022–2026 from mfapi.in"),
        ("90 rows", "aum_by_fund_house.csv — Quarterly AUM per AMC"),
        ("48 rows", "monthly_sip_inflows.csv — Industry SIP trend data"),
        ("32,778 rows", "investor_transactions.csv — SIP/Lumpsum/Redemption records"),
        ("8,050 rows", "benchmark_indices.csv — Nifty 50, 100, BSE SmallCap daily values"),
        ("322 rows", "portfolio_holdings.csv — Top stock holdings per equity fund"),
    ]
    for i, (count, desc) in enumerate(datasets):
        add_text(slide, count, 0.5, 1.2 + i*0.7, 1.8, 0.5, fontsize=13, bold=True, color=PURPLE)
        add_text(slide, desc, 2.5, 1.2 + i*0.7, 10, 0.5, fontsize=11, color=DARK)

    # Slide 4 — Architecture
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs, DARK)
    add_text(slide, "System Architecture", 0.5, 0.3, 12, 0.8, fontsize=28, bold=True, color=WHITE)
    layers = [
        ("EXTRACT", "mfapi.in REST API  +  AMFI CSV datasets", PURPLE),
        ("TRANSFORM", "Pandas cleaning  •  Date fixing  •  Forward-fill  •  Validation", ORANGE),
        ("LOAD", "SQLite star schema  •  dim_fund  •  fact_nav  •  fact_transactions", PURPLE),
        ("ANALYSE", "CAGR  •  Sharpe  •  Alpha/Beta  •  VaR  •  Cohort Analysis", ORANGE),
        ("VISUALISE", "Power BI 4-page Dashboard  •  9 Matplotlib/Plotly Charts", PURPLE),
    ]
    for i, (layer, desc, color) in enumerate(layers):
        add_text(slide, layer, 0.5, 1.3 + i*1.1, 2, 0.5, fontsize=14, bold=True, color=color)
        add_text(slide, desc, 2.8, 1.3 + i*1.1, 10, 0.5, fontsize=12, color=WHITE)

    # Slide 5 — EDA Highlights 1
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs, WHITE)
    add_text(slide, "EDA Highlights — Industry Trends", 0.5, 0.3, 12, 0.8, fontsize=28, bold=True, color=PURPLE)
    findings = [
        "📈  SIP inflows tripled from ₹11,517 Cr (Jan 2022) to ₹31,002 Cr (Dec 2025) — all-time high",
        "🏦  SBI Mutual Fund dominates AUM at ₹12.5L Cr — nearly 2x ICICI Prudential",
        "👥  Industry folios doubled from 13.26 Cr to 26.12 Cr between 2022 and 2025",
        "💧  Liquid funds attract ₹33,000–42,000 Cr monthly — largest category by far",
        "🏙️  T30 cities drive 66.3% of transactions vs 33.7% from B30 cities",
    ]
    for i, f in enumerate(findings):
        add_text(slide, f, 0.5, 1.3 + i*1.0, 12, 0.8, fontsize=13, color=DARK)
    
    # Add NAV chart if exists
    chart_path = CHARTS_DIR / '01_nav_trend.png'
    if chart_path.exists():
        slide.shapes.add_picture(str(chart_path), Inches(7.5), Inches(1.5), Inches(5.5), Inches(3.5))

    # Slide 6 — EDA Highlights 2
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs, WHITE)
    add_text(slide, "EDA Highlights — Investor Demographics", 0.5, 0.3, 12, 0.8, fontsize=28, bold=True, color=PURPLE)
    findings2 = [
        "👶  26-35 age group dominates — 41.1% of all investors",
        "👨  Gender gap persists — 66.5% male vs 33.5% female investors",
        "🏢  Banking (19.2%), IT (13.4%), Pharma (12.0%) = 44.6% of equity portfolios",
        "💰  Average SIP amounts are consistent across age groups (~₹100K)",
        "📊  Small Cap funds delivered highest 1yr returns — DSP Small Cap at 64.88%",
    ]
    for i, f in enumerate(findings2):
        add_text(slide, f, 0.5, 1.3 + i*1.0, 7, 0.8, fontsize=13, color=DARK)
    
    chart_path2 = CHARTS_DIR / '05_demographics.png'
    if chart_path2.exists():
        slide.shapes.add_picture(str(chart_path2), Inches(7.5), Inches(1.5), Inches(5.5), Inches(3.5))

    # Slide 7 — Performance Metrics
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs, WHITE)
    add_text(slide, "Fund Performance Metrics", 0.5, 0.3, 12, 0.8, fontsize=28, bold=True, color=PURPLE)
    
    scorecard = pd.read_csv(PROCESSED_DIR / 'fund_scorecard.csv')
    top5 = scorecard.nlargest(5, 'score')[['scheme_name', 'score', 'cagr_3yr_pct', 'sharpe_ratio']]
    
    headers = ['Fund', 'Score', '3yr CAGR%', 'Sharpe']
    add_text(slide, ' | '.join(headers), 0.5, 1.2, 12, 0.4, fontsize=11, bold=True, color=PURPLE)
    for i, (_, row) in enumerate(top5.iterrows()):
        line = f"{row['scheme_name'][:45]}  |  {round(row['score'],1)}  |  {round(row['cagr_3yr_pct'],2)}%  |  {round(row['sharpe_ratio'],3)}"
        add_text(slide, line, 0.5, 1.7 + i*0.6, 12, 0.5, fontsize=10, color=DARK)
    
    metrics = [
        "Sharpe Ratio = (Rp - Rf) / σ × √252  |  Rf = 6.5% RBI repo rate",
        "CAGR = (NAV_end / NAV_start)^(1/n) - 1",
        "Max Drawdown = min(NAV / running_max - 1)",
        "Fund Score = 30%×CAGR + 25%×Sharpe + 20%×Alpha + 15%×Expense + 10%×MaxDD",
    ]
    for i, m in enumerate(metrics):
        add_text(slide, m, 0.5, 4.5 + i*0.5, 12, 0.4, fontsize=10, color=ORANGE)

    # Slide 8 — Risk Metrics
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs, DARK)
    add_text(slide, "Advanced Risk Metrics", 0.5, 0.3, 12, 0.8, fontsize=28, bold=True, color=WHITE)
    risk_metrics = [
        ("VaR (95%)", "ABSL Small Cap: -2.39% daily VaR — worst 5% of days exceed this loss"),
        ("CVaR", "Average loss on bad days: -3.03% for Small Cap vs -1.2% for Large Cap"),
        ("Rolling Sharpe", "Sharpe fluctuates -4 to +6 — no fund consistently outperforms"),
        ("SIP Continuity", "97.8% of investors show irregular SIP gaps (avg 64.9 days vs 30 expected)"),
        ("Sector HHI", "Axis Bluechip most concentrated (3000) vs UTI Mid Cap most diversified (1300)"),
    ]
    for i, (metric, desc) in enumerate(risk_metrics):
        add_text(slide, metric, 0.5, 1.3 + i*1.0, 2.5, 0.5, fontsize=13, bold=True, color=ORANGE)
        add_text(slide, desc, 3.2, 1.3 + i*1.0, 9.5, 0.5, fontsize=12, color=WHITE)

    # Slide 9 — Dashboard Page 1 & 2
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs, WHITE)
    add_text(slide, "Power BI Dashboard", 0.5, 0.3, 12, 0.8, fontsize=28, bold=True, color=PURPLE)
    add_text(slide, "Page 1 — Industry Overview", 0.5, 1.0, 6, 0.4, fontsize=13, bold=True, color=ORANGE)
    add_text(slide, "Page 2 — Fund Performance", 6.8, 1.0, 6, 0.4, fontsize=13, bold=True, color=ORANGE)
    
    p1 = BASE_DIR / 'dashboard' / 'page1_industry_overview.png'
    p2 = BASE_DIR / 'dashboard' / 'page2_fund_performance.png'
    if p1.exists():
        slide.shapes.add_picture(str(p1), Inches(0.5), Inches(1.5), Inches(6), Inches(4))
    if p2.exists():
        slide.shapes.add_picture(str(p2), Inches(6.8), Inches(1.5), Inches(6), Inches(4))

    # Slide 10 — Dashboard Page 3 & 4
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs, WHITE)
    add_text(slide, "Power BI Dashboard", 0.5, 0.3, 12, 0.8, fontsize=28, bold=True, color=PURPLE)
    add_text(slide, "Page 3 — Investor Analytics", 0.5, 1.0, 6, 0.4, fontsize=13, bold=True, color=ORANGE)
    add_text(slide, "Page 4 — SIP & Market Trends", 6.8, 1.0, 6, 0.4, fontsize=13, bold=True, color=ORANGE)
    
    p3 = BASE_DIR / 'dashboard' / 'page3_investor_analytics.png'
    p4 = BASE_DIR / 'dashboard' / 'page4_sip_market_trends.png'
    if p3.exists():
        slide.shapes.add_picture(str(p3), Inches(0.5), Inches(1.5), Inches(6), Inches(4))
    if p4.exists():
        slide.shapes.add_picture(str(p4), Inches(6.8), Inches(1.5), Inches(6), Inches(4))

    # Slide 11 — Key Findings
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs, WHITE)
    add_text(slide, "Key Findings & Recommendations", 0.5, 0.3, 12, 0.8, fontsize=28, bold=True, color=PURPLE)
    findings_recs = [
        "🏆  ICICI Pru Midcap tops Fund Scorecard at 85.1 — best risk-return balance",
        "📉  Small Cap funds: highest returns (64.88%) but also highest VaR (-2.39%)",
        "💡  Moderate risk investors: Mirae Asset Large Cap offers best Sharpe (1.068)",
        "⚠️   SIP continuity critical — 97.8% investors show irregular payment patterns",
        "🌍  B30 cities at 33.7% — significant untapped growth opportunity",
        "👩  Female investor gap (33.5%) — targeted outreach programs recommended",
        "🏦  Debt funds underperform risk-free rate — negative Sharpe across all debt funds",
    ]
    for i, f in enumerate(findings_recs):
        add_text(slide, f, 0.5, 1.2 + i*0.75, 12, 0.6, fontsize=13, color=DARK)

    # Slide 12 — Thank You
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs, DARK)
    add_text(slide, "Thank You", 1, 2.0, 11, 1.2, fontsize=48, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(slide, "Mutual Fund Analytics Platform — Bluestock Fintech Internship", 1, 3.5, 11, 0.6, fontsize=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "GitHub Repository  |  July 2026", 1, 4.5, 11, 0.5, fontsize=12, color=ORANGE, align=PP_ALIGN.CENTER)

    prs.save(str(REPORTS_DIR / 'Presentation.pptx'))
    print("✅ Presentation.pptx generated!")

if __name__ == "__main__":
    generate_pptx()